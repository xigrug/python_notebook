# -*- coding:utf-8 -*-

"""
See http://pbpython.com/creating-powerpoint.html for details on this script
Requires https://python-pptx.readthedocs.org/en/latest/index.html
Example program showing how to read in Excel, process with pandas and
output to a PowerPoint file.
"""

from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches
import argparse
import pandas as pd
import numpy as np
from datetime import date, timedelta
import matplotlib.pyplot as plt
import seaborn as sns
####email
import smtplib    #发邮件
from email.mime.text import MIMEText   #邮件文字
from email.mime.multipart import MIMEMultipart   #邮件附件
import os     #系统命令
#import sys
#reload(sys)
#sys.setdefaultencoding('utf-8')

def df_to_table(slide, df, left, top, width, height, colnames=None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.
    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data
    Optional arguments:
     - colnames
     https://github.com/robintw/PandasToPowerpoint/blob/master/PandasToPowerpoint.py
     """
    rows, cols = df.shape
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            col_name = " ".join(col_name)
        res.table.cell(0, col_index).text = col_name

    m = df.as_matrix()

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text


def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Create ppt report')
    parser.add_argument('infile',
                        type=argparse.FileType('r'),
                        help='Powerpoint file used as the template')
    parser.add_argument('report',
                        type=argparse.FileType('r'),
                        help='Excel file containing the raw report data')
    parser.add_argument('outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint report file')
    return parser.parse_args()


def create_pivot(df, index_list=["Manager", "Rep", "Product"],
                 value_list=["Price", "Quantity"]):
    """
    Take a DataFrame and create a pivot table
    Return it as a DataFrame pivot table
    """
    table = pd.pivot_table(df, index=index_list,
                           values=value_list,
                           aggfunc=[np.sum, np.mean], fill_value=0)
    return table


def create_chart(df, filename):
    """ Create a simple bar chart saved to the filename based on the dataframe
    passed to the function
    """
    df['total'] = df['Quantity'] * df['Price']
    final_plot = df.groupby('Name')['total'].sum().order().plot(kind='barh')
    fig = final_plot.get_figure()
    fig.set_size_inches(6, 4.5)
    fig.savefig(filename, bbox_inches='tight', dpi=600)
def getYesterday(): 
    today=date.today() 
    oneday=timedelta(days=1) 
    yesterday=today-oneday  
    return yesterday
def getTomorrow():
    today=date.today()
    oneday=timedelta(days=1)
    tomorrow=today+oneday
    return tomorrow
def getday(tday):
    today=date.today()
    oneday=timedelta(days=tday)
    day=today+oneday
    return day
def creat_slide2fig8(prs,tx1,fig1,fig2,fig3,fig4,fig5,fig6,fig7,fig8):
    graph_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = tx1
    subtitle = slide.placeholders[1]
    subtitle.text = " {:%m-%d-%Y}".format(getday(1))
    subtitle = slide.placeholders[3]
    subtitle.text = " {:%m-%d-%Y}".format(getday(2))
    subtitle = slide.placeholders[22]
    subtitle.text = " {:%m-%d-%Y}".format(getday(3))
    subtitle = slide.placeholders[23]
    subtitle.text = " {:%m-%d-%Y}".format(getday(4))
    subtitle = slide.placeholders[26]
    subtitle.text = " {:%m-%d-%Y}".format(getday(5))
    subtitle = slide.placeholders[25]
    subtitle.text = " {:%m-%d-%Y}".format(getday(6))
    subtitle = slide.placeholders[24]
    subtitle.text = " {:%m-%d-%Y}".format(getday(7))
    subtitle = slide.placeholders[27]
    subtitle.text = " {:%m-%d-%Y}".format(getday(8))
    placeholder = slide.placeholders[16]
    pic = placeholder.insert_picture(fig1)
    placeholder = slide.placeholders[15]
    pic = placeholder.insert_picture(fig2)
    placeholder = slide.placeholders[17]
    pic = placeholder.insert_picture(fig3)
    placeholder = slide.placeholders[13]
    pic = placeholder.insert_picture(fig4)
    placeholder = slide.placeholders[20]
    pic = placeholder.insert_picture(fig5)
    placeholder = slide.placeholders[19]
    pic = placeholder.insert_picture(fig6)
    placeholder = slide.placeholders[21]
    pic = placeholder.insert_picture(fig7)
    placeholder = slide.placeholders[18]
    pic = placeholder.insert_picture(fig8)
    
def creat_slide2fig2(prs,tx1,fig1,fig2):
    graph_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = tx1
    subtitle = slide.placeholders[1]
    subtitle.text = "{:%Y-%m-%d}".format(getday(1))+" / {:%Y-%m-%d}".format(getday(8))+' 08'
    placeholder = slide.placeholders[13]
    pic = placeholder.insert_picture(fig1)
    subtitle = slide.placeholders[3]
    subtitle.text = "{:%Y-%m-%d}".format(getday(1))+" / {:%Y-%m-%d}".format(getday(8))+' 20'
    placeholder = slide.placeholders[14]
    pic = placeholder.insert_picture(fig2)
def creat_slide2fig22(prs,tx1,fig1,fig2):
    graph_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = tx1
    subtitle = slide.placeholders[1]
    subtitle.text = "{:%Y-%m-%d}".format(getday(1))+" / {:%Y-%m-%d}".format(getday(8))+' 14'
    placeholder = slide.placeholders[13]
    pic = placeholder.insert_picture(fig1)
    subtitle = slide.placeholders[3]
    subtitle.text = "{:%Y-%m-%d}".format(getday(2))+" / {:%Y-%m-%d}".format(getday(9))+' 02'
    placeholder = slide.placeholders[14]
    pic = placeholder.insert_picture(fig2)
def creat_slide2fig1(prs,tx1,fig1):
    graph_slide_layout = prs.slide_layouts[10]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = tx1
    placeholder = slide.placeholders[1]
    pic = placeholder.insert_picture(fig1)

def create_ppt(input, output, chart):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    df = pd.read_csv(chart,header=None,names=["name"])
    print(df.name[0])
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "WRF 气象场预报" 
    subtitle.text = " {:%m-%d-%Y}".format(getTomorrow())
    # Create the summary graph
    creat_slide2fig2(prs,"500hPa 天气形势",df.name[65],df.name[66])
    creat_slide2fig2(prs,"地面 天气形势",df.name[67],df.name[68])
    creat_slide2fig1(prs,"d04 风速风向各站逐小时变化",df.name[51])
    creat_slide2fig1(prs,"d04 降水各站逐小时变化",df.name[52])
    creat_slide2fig8(prs,"d03 日均 2m温度",df.name[0],df.name[1],df.name[2],df.name[3],df.name[4],df.name[5],df.name[6],df.name[7])
    creat_slide2fig8(prs,"d04 日均 2m温度",df.name[8],df.name[9],df.name[10],df.name[11],df.name[12],\
                     df.name[13],df.name[14],df.name[15])
    creat_slide2fig1(prs,"d04 温度各站逐小时变化",df.name[48])
    creat_slide2fig8(prs,"d03 日均 2m湿度",df.name[16],df.name[17],df.name[18],df.name[19],df.name[20],\
                     df.name[21],df.name[22],df.name[23])
    creat_slide2fig8(prs,"d04 日均 2m湿度",df.name[24],df.name[25],df.name[26],df.name[27],df.name[28],\
                     df.name[29],df.name[30],df.name[31])
    creat_slide2fig1(prs,"d04 湿度各站逐小时变化",df.name[49])
    creat_slide2fig1(prs,"d04 总云量各站逐小时变化",df.name[73])
    creat_slide2fig1(prs,"d04 低云量各站逐小时变化",df.name[74])
    creat_slide2fig1(prs,"d04 短波辐射各站逐小时变化",df.name[75])
    creat_slide2fig22(prs,"垂直廓线",df.name[70],df.name[72])
    # Save PPT
    prs.save(output)
def send_email(output):
    msg=MIMEMultipart()
    msg['from'] = 'xigrug@126.com'+'<'+'xigrug@126.com'+'>'
    msg['to'] = 'liuchongcn@126.com'
    msg['subject'] = '人生苦短weather'
    txt=MIMEText("这是正文内容，附件附件附件，请查收",'plain','utf-8')
    msg.attach(txt)
    #构造附件1
    att1=MIMEText(open(output, 'rb').read(), 'base64', 'utf-8')
    att1["Content-Type"] = 'application/octet-stream'
    att1["Content-Disposition"] = 'attachment; filename=%s'%(output)#这里的filename可以任意写，写什么名字，邮件中显示什么名字
    msg.attach(att1)

    try:
        server = smtplib.SMTP()        #构造邮件传输服务
        server.connect('smtp.126.com')   #连接qq邮件host
        server.starttls()
        server.login('xigrug@126.com','GRUG12350')#前为用户名，后为密码
        server.sendmail(msg['from'], msg['to'],msg.as_string())   #从哪到哪，传邮件
        server.quit()  #退出邮件服务
        print('发送成功')
    except Exception as e:
        print(str(e))
    exit()  ####此处将exit()注释掉则执行立即关机，否则运行止于此


if __name__ == "__main__":
    args = parse_args()
    create_ppt(args.infile.name, args.outfile.name, args.report.name)
    #send_email(args.outfile.name)